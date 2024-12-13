import json
from pathlib import Path

import logfire
import pandas as pd
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
from docx import Document
from httpx import Client
from pptx import Presentation

from app.models.file import FileType


def read_content(file_path: Path, file_type: FileType) -> str | None:
    """
    Read content from file based on file type
    Args:
        file_path: Path to the file
        file_type: Type of the file (PDF, DOCX, etc.)
    Returns:
        str | None: Content of the file or None if unsupported/error
    """
    try:
        content = ""

        if file_type == FileType.pdf:
            reader = PdfReader(file_path)
            for page in reader.pages:
                content += page.extract_text()

        elif file_type == FileType.docx:
            doc = Document(file_path)
            content = '\n'.join(para.text for para in doc.paragraphs)

        elif file_type == FileType.txt:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

        elif file_type == FileType.pptx:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + '\n'

        elif file_type == FileType.md:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

        elif file_type in [FileType.csv, FileType.xlsx]:
            df = pd.read_csv(file_path) if file_type == FileType.csv else pd.read_excel(file_path)
            content = df.to_string()

        elif file_type == FileType.html:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, "html.parser")
                content = soup.get_text()

        elif file_type == FileType.json:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                content = json.dumps(data, indent=2)

        else:
            logfire.error(f"Unsupported file type: {file_type}")
            return None

        logfire.info(f"Content extracted from {file_path} of type {file_type} \n"
                    f"{content[:100]}")
        return content

    except Exception as e:
        logfire.error(f"Error reading file content: {str(e)}")
        raise


def read_url_content(url: str) -> str:
    """
    Read content from URL
    Args:
        url: URL to read content from
    Returns:
        str: Content from the URL
    """
    try:
        client = Client()
        response = client.get(url)
        soup = BeautifulSoup(response.text, "html.parser")
        content = soup.get_text()
        logfire.info(f"Content extracted from {url} \n"
                    f"{content[:100]}")
        return content
    except Exception as e:
        logfire.error(f"Error reading URL content: {str(e)}")
        raise
